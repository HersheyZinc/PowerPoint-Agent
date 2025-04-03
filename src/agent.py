from pptx import Presentation
from .utils import fromEmus, ppt_to_pdf, pdf_to_img
from .ppt_reader import get_slide_content, get_ppt_content
from .ppt_writer import insert_slide, delete_all_shapes
from .openai import query, query_tools
from . import apis, prompts
from PIL import Image
import json, os, tempfile, shutil, datetime
from concurrent.futures import ThreadPoolExecutor


# TODO
# Ingest existing presentations -> Iterate through each slide and generate a summary


class AgentPPT():
    def __init__(self, model="gpt-4o", src_path="", dst_path="working.pptx"):
        self.ppt = None
        self.ppt_path = dst_path
        self.slide_idx = 0
        self.verbose=True
        self.threading=True

        # LLM args
        self.model = model
        self.model_temp = 0

        self.chat_history = []
        
        self.logger = []

        self.log(f"AgentPPT instance created with model: {model}, src_path: src_path, dst_path: {dst_path}")
        self.new_ppt(src_path)


    def new_ppt(self, file_path=""):
        if file_path:
            self.ppt = Presentation(file_path)
        else:
            self.ppt = Presentation()
            # self.insert_slide()

        self.slide_idx = 0
        self.clear_chat_history()
        self.log("New presentation created")
        

    def save_ppt(self):
        # Write content to pptx file
        self.ppt.save(self.ppt_path)
        self.log(f"Presentation saved to {self.ppt_path}")


    def clear_chat_history(self):
        system_prompt = "You are an expert Powerpoint slide designer. Use the tools provided to fulfil the user's request."
        self.chat_history = []
        self.log("Cleared chat history.")

    
    def log(self, log_str):
        date_str = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        formatted_str = f"{date_str} - {log_str}"
        self.logger.append(formatted_str)
        if self.verbose:
            print(formatted_str, "\n")


    def print_chat_history(self):
        for msg in self.chat_history:
            print(f"{msg["role"]}:\n{msg["content"]}\n{'--------------'*10}")


    def print_ppt(self):
        print(get_ppt_content(self.ppt))
    
        
    def render(self):
        if len(self.ppt.slides) == 0:
            white_image = Image.new("RGB", (800, 600), "white")
            return [white_image]

        temp_dir = tempfile.mkdtemp()

        temp_ppt_path = os.path.join(temp_dir, "temp_presentation.pptx")
        self.ppt.save(temp_ppt_path)
        pdf_path = ppt_to_pdf(temp_ppt_path, temp_dir)
        slide_images = pdf_to_img(pdf_path)

        shutil.rmtree(temp_dir)
        self.log("Rendered slides to images.")
        return slide_images


    def generate_module(self, prompt):
        self.log(f"Generating PowerPoint presentation with prompt: {prompt}")
        self.chat_history.append({"role":"user", "content": prompt})
        messages = [{"role":"system", "content":prompts.generate_prompt}, {"role":"user", "content": prompt}]
        toolkit = [a.get_openai_args() for a in apis.OUTLINES]
        response = query(messages,  model=self.model)
        
        for line in response.split("\n"):
            if line.startswith("Title:"):
                title = line[6:].strip()
                self.insert_slide()
                slide = self.ppt.slides[-1]
                slide.shapes[0].text = title
            elif line.startswith("-"):
                text = line[1:].strip()
                slide.shapes[1].text += text + "\n"

        # output_str = f"{len(tool_calls)} slides inserted."
        # self.log(output_str)
        # for slide_idx, tool_call in enumerate(tool_calls):
        #     fn_args = json.loads(tool_call.function.arguments)
        #     title = fn_args["title"]
        #     content = fn_args["text"]
        #     self.insert_slide()
        #     output_str + "\n" + self.action_module(f"Change the title to {title}. Change the content to {content}", slide_idx)


        messages = [{"role":"system", "content": prompts.user_response_prompt}, {"role":"system", "content":response}]
        agent_response = query(messages, model=self.model, temperature=0.2)
        self.chat_history.append({"role":"assistant", "content":agent_response})
        self.log("Agent response: \n" + agent_response)
        return agent_response


    def plan_module(self, prompt):
        self.log(f"Calling PLAN module with prompt: {prompt}")

        self.chat_history.append({"role":"user", "content": prompt})
        ppt_content = get_ppt_content(self.ppt)
        
        # if len(self.chat_history) == 1:
        #     messages = [{"role":"system", "content":prompts.enhance_prompt}, {"role":"user", "content": prompt}]
        #     enhanced_prompt = query(messages, model="gpt-4o-mini", temperature=0.2)
        #     self.log("Enhanced orginal prompt:\n"+enhanced_prompt)

        #     messages = [{"role":"system", "content":prompts.plan_prompt}, {"role":"system", "content":ppt_content},{"role":"user", "content":enhanced_prompt}]
        # else:
        messages = [{"role":"system", "content":prompts.plan_prompt}, {"role":"system", "content":ppt_content}] + self.chat_history[-5:] + [{"role":"system", "content":f"The user is currently looking at slide {self.slide_idx}"}]

        toolkit = [a.get_openai_args() for a in apis.PLANS]
        _, tool_calls = query_tools(messages, toolkit, model=self.model, max_tokens=16000)

        output_str = ""

        for tool_call in tool_calls:
            fn_name = tool_call.function.name
            fn_args = json.loads(tool_call.function.arguments)
            if fn_name == 'insert_slide':
                template_request = fn_args["slide_template"]
                response = insert_slide(self.ppt, template_request, self.model)
                output_str += response + "\n"
                slide_idx = len(self.ppt.slides)-1
                self.log(response)

            elif fn_name == "redo_slide":
                slide_idx = fn_args['slide_index']
                response = delete_all_shapes(self.ppt, slide_idx)
                output_str += response + "\n"
                self.log(response)

            else:
                slide_idx = fn_args['slide_index']


            if "instructions" in fn_args:
                output_str += f"Modified slide {slide_idx} with instructions: {fn_args['instructions']}\n"
                output_str += self.action_module("Modify the slide according to the following instructions: " + fn_args["instructions"], slide_idx)


            
        messages = [{"role":"system", "content": prompts.user_response_prompt}, {"role":"system", "content":output_str}]
        agent_response = query(messages, model=self.model, temperature=0.2)
        self.chat_history.append({"role":"assistant", "content":agent_response})
        self.log("Agent response: \n" + agent_response)
        return agent_response


    def action_module(self, prompt, slide_idx):
        self.log(f"Calling ACTION module on slide {slide_idx} with prompt: {prompt}")

        def execute_tool_call(tool_call):
            fn_name = tool_call.function.name
            fn_args = json.loads(tool_call.function.arguments)
            for api in apis.ACTIONS:
                if api.name == fn_name:
                    try:
                        r = api.run(self.ppt, slide_idx, fn_args, self.model)
                        return f"API - {fn_name} | Status - SUCCESS | {r} | Arguments - {fn_args}"
                    except Exception as e:
                        return f"API - {fn_name} | Status - FAILED | parameters - {fn_args} | {e}"
                    


        slide_content = get_slide_content(self.ppt, slide_idx)
        slide = self.ppt.slides[slide_idx]
        messages = [{"role":"system", "content":prompts.action_prompt}, {"role":"system", "content":slide_content}, {"role":"user", "content":prompt}]
        toolkit = [a.get_openai_args() for a in apis.ACTIONS]
        _, tool_calls = query_tools(messages, toolkit, model=self.model)

        if self.threading:
            with ThreadPoolExecutor() as executor:
                results = executor.map(execute_tool_call, tool_calls)
                output_str = "\n".join(results)
        else:
            output_str = ""
            for tool_call in tool_calls:
                output_str += execute_tool_call(tool_call) + "\n"



        self.log(output_str)
        return output_str
    


    def insert_slide(self, layout_idx=1, name="", summary="", script=""):
        slide_layout = self.ppt.slide_layouts[layout_idx]
        slide = self.ppt.slides.add_slide(slide_layout)

        slide.name = name
        slide.description = f"An empty slide with layout - {slide_layout.name}"
        slide.summary = summary
        slide.script = script

        return "Slide inserted."
        
