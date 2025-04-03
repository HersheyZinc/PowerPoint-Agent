import collections 
import collections.abc
import pptx.parts.image
import pptx.enum.shapes as shapes
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.dml.color import RGBColor
from .utils import fromEmus, fromPts
import json

global slides
global ppt
global global_args
SCALE = 1000
slide_height = 0
slide_width = 0
shape_list = ['PLACEHOLDER', 'PICTURE', 'CHART', 'TABLE', 'TEXT_BOX', 'AUTO_SHAPE']
slides = None


def get_fill_color(shape):
    if shape.fill.type == 1:  # Solid fill
        color = shape.fill.fore_color
        if hasattr(color, "rgb"):
            return color.rgb
    return None

def str2json(str):
    json = {}
    for pair in [s.strip() for s in str.split(",")]:
        key, value = pair.split('=')
        json[key] = value

    return json


class BasicShape:
    def __init__(self, shape):
        self.shape_type = shape.shape_type
        self.height = shape.height
        self.width = shape.width
        self.left = shape.left
        self.top = shape.top
        self.name = shape.name
        self.shape_id = shape.shape_id
    
    @property
    def text_info(self):
        pass

    @property
    def space_info(self):
        return f"left={fromEmus(self.left)}, top={fromEmus(self.top)}"
    
    @property
    def size_info(self):
        return f"height={fromEmus(self.height)}, width={fromEmus(self.width)}"

    @property
    def style_info(self):
        pass

    @property
    def description(self):
        # return f"[{self.name.split(' ')[0]}]\n" 
        try:
            return f"{self.name.split(' ')[0]}"
        except:
            return f"{str(self.shape_type).split(' ')[0].strip()}" 

    def __repr__(self):
        s = ""
        s += self.description
        s += self.size_info
        if self.text_info is not None:
            s += self.text_info
        s += self.space_info
        if self.style_info is not None:
            s += self.style_info
        return s


class Picture(BasicShape):
    def __init__(self, shape, id=None):
        super().__init__(shape)
        self.image = shape.image
        self.rotation = int(shape.rotation)
        self.id = id
    
    @property
    def style_info(self):
        return f"rotation={self.rotation}"
    
    @property
    def description(self):
        return "Picture"

class Table(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.table = shape.table
        self.rows = shape.table.rows
        self.columns = shape.table.columns

    @property
    def text_info(self):
        s = "Data:\n"
        for row in self.rows:
            s += "|"
            for col in row.cells:
                s += f"{col.text}|"
            s += "\n"
        return s
    
    @property
    def description(self):
        return f"Table with {len(self.rows)} rows and {len(self.columns)} columns" 

class Chart(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.chart = shape.chart
        self.title = shape.chart.chart_title.text_frame.text
        self.chart_type = str(shape.chart.chart_type).split(' ')[0]
    
    @property
    def text_info(self):
        s = ""
        if self.title:
            s += f"Title: {self.title}\n"
        s += f"Chart Type: {self.chart_type}\n"
        s += "Data:\n"
        try:
            for series in self.chart.series:
                s += f"{series.name}: "
                for value in series.values:
                    s += f"{value}, "
                s += "\n"
        except:
            pass
        return s
    
    @property
    def style_info(self):
        return ""
    
    @property
    def description(self):
        return "Chart"


class Textbox(BasicShape):
    def __init__(self, shape, id=None):
        super().__init__(shape)
        self.text = shape.text_frame.text
        self.paragraphs = shape.text_frame.paragraphs
        try:
            self.font = self.paragraphs[0].runs[0].font
        except:
            self.font = self.paragraphs[0].font
        self.bold = self.font.bold
        self.italic = self.font.italic
        self.underline = self.font.underline
        self.size = self.font.size if self.font.size is not None else self.paragraphs[0].font.size
        self.size = fromPts(self.size)
        try:
            self.color = self.font.color.rgb 
        except:
            self.color = None
        self.fill = get_fill_color(shape)
        self.font_name = self.font.name
        self.line_spacing = self.paragraphs[0].line_spacing
        self.align = self.paragraphs[0].alignment
        self.id=id
    
    @property
    def text_info(self):
        return f"{self.text}"
    
    @property
    def style_info(self):
        return f'bold={self.bold}, italic={self.italic}, underline={self.underline}, font size={self.size}, color={self.color}, fill={self.fill}, font style={self.font_name}, line_space={self.line_spacing}, align={self.align}'

    @property
    def description(self):
        return "TextBox"
    
class Placeholder(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.fill = get_fill_color(shape)
        self.text = shape.text_frame.text
        if shape.has_text_frame:
            textframe = shape.text_frame
            try:
                font = shape.text_frame.paragraphs[0].runs[0].font
            except:
                font = shape.text_frame.paragraphs[0].font
            self.bold = font.bold
            self.italic = font.italic
            self.underline = font.underline
            self.size = fromPts(font.size) if font.size else None
            try:
                self.color = font.color.rgb 
            except:
                self.color = None
            self.font_name = font.name
            self.line_spacing = textframe.paragraphs[0].line_spacing
            self.align = textframe.paragraphs[0].alignment
    
    @property
    def text_info(self):
        if self.text is not None:
            return f"{self.text}"
        else:
            return ""
    
    @property
    def style_info(self):
        return f'bold={self.bold}, italic={self.italic}, underline={self.underline}, font size={self.size}, color={self.color}, fill={self.fill}, font style={self.font_name}, line_space={self.line_spacing}, align={self.align}\n'



class AutoShape(BasicShape):
    def __init__(self, shape):
        super().__init__(shape)
        self.text = shape.text_frame.text
        self.fill = get_fill_color(shape)
    
    @property
    def text_info(self):
        return f"{self.text}"
    
    @property
    def style_info(self):
        return f"fill={self.fill}"
        # return ""

    
def hasshape(shape_str, shape_list):
    for shape in shape_list:
        if shape in shape_str:
            return True
    return False


def get_slide_content(ppt, slide_idx, return_json=False):
    slide = ppt.slides[slide_idx]
    output = {"slide_idx":slide_idx, "slide_height":fromEmus(ppt.slide_height), "slide_width":fromEmus(ppt.slide_width),"slide_background":get_fill_color(slide.background)}
    if slide.notes_slide.notes_text_frame.text:
            output["slide_notes"] = slide.notes_slide.notes_text_frame.text
    
    elements = []
    for shape_idx, shape in enumerate(slide.shapes):
        if 'PLACEHOLDER' in str(shape.shape_type):
            shape = Placeholder(shape)
        elif 'PICTURE' in str(shape.shape_type):
            shape = Picture(shape)
        elif 'CHART' in str(shape.shape_type):
            shape = Chart(shape)
        elif 'TABLE' in str(shape.shape_type):
            shape = Table(shape)
        elif 'TEXT_BOX' in str(shape.shape_type):
            shape = Textbox(shape)
        elif 'AUTO_SHAPE' in str(shape.shape_type):
            shape = AutoShape(shape)
        else:
            shape = AutoShape(shape)

        element_info = {"index":shape_idx, "type":shape.description, "size":str2json(shape.size_info)}

        if shape.text_info:
            element_info["text"] = shape.text_info
        if shape.style_info:
            element_info["text_style"] = str2json(shape.style_info)
        if shape.space_info:
            element_info["position"] = str2json(shape.space_info)

        elements.append(element_info)

    output["shapes"] = elements

    if return_json:
        return output
    return json.dumps(output, indent=2)


def get_ppt_content(ppt):
    ppt_content = []
    for slide_idx, slide in enumerate(ppt.slides):
        ppt_content.append(get_slide_content(ppt, slide_idx, return_json=True))
    if not ppt_content:
        return "There are no slides."
    return json.dumps(ppt_content, indent=2)


def get_shape_content(slide, shape_idx, return_json=False):
    shape = slide.shapes[shape_idx]
    if 'PLACEHOLDER' in str(shape.shape_type):
        shape = Placeholder(shape)
    elif 'PICTURE' in str(shape.shape_type):
        shape = Picture(shape)
    elif 'CHART' in str(shape.shape_type):
        shape = Chart(shape)
    elif 'TABLE' in str(shape.shape_type):
        shape = Table(shape)
    elif 'TEXT_BOX' in str(shape.shape_type):
        shape = Textbox(shape)
    elif 'AUTO_SHAPE' in str(shape.shape_type):
        shape = AutoShape(shape)
    else:
        shape = AutoShape(shape)

    shape_content = {"index":shape_idx, "type":shape.description, "size":str2json(shape.size_info)}

    if shape.text_info:
        shape_content["text"] = shape.text_info
    if shape.style_info:
        shape_content["text_style"] = str2json(shape.style_info)
    if shape.space_info:
        shape_content["position"] = str2json(shape.space_info)

    if return_json:
        return shape_content
    return json.dumps(shape_content, indent=2)





