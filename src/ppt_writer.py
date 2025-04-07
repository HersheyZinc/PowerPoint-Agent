from pptx.dml.color import RGBColor
from .utils import validate_hex, toEmus
from pptx.util import Pt
from pptx.parts.image import Image
from .ppt_reader import get_shape_content
from pptx.enum.shapes import MSO_SHAPE
from .openai import generate_image, query
from .prompts import *
from io import BytesIO
import json
SHAPE_DICT = {getattr(MSO_SHAPE, attr):attr for attr in dir(MSO_SHAPE) if attr.isupper()}
SHAPE_DICT = json.dumps(dict(sorted(SHAPE_DICT.items())))


def set_shape_properties(shape, parameters):
    if "top" in parameters:
        shape.top = toEmus(parameters["top"])
    if "left" in parameters:
        shape.left = toEmus(parameters["left"])
    if "height" in parameters:
        shape.height = toEmus(parameters["height"])
    if "width" in parameters:
        shape.width = toEmus(parameters["width"])
    if "fill_color" in parameters:
        fill_color = parameters["fill_color"]
        if fill_color =="transparent":
            shape.fill.background()
        elif validate_hex(fill_color):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color.replace("#",""))
            

    if "has_border" in parameters:
        if parameters["has_border"]:
            shape.line.fill.solid()
        else:
            shape.line.fill.background()
    if "border_width" in parameters:
        shape.line.width = Pt(parameters["border_width"])
    if "border_color" in parameters and validate_hex(parameters["border_color"]):
        shape.line.color.rgb = RGBColor.from_string(parameters["border_color"].replace("#",""))


    if "text" in parameters:
        shape.text = parameters["text"]
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in [run for run in paragraph.runs] + [paragraph]:
                
                if "font_color" in parameters and validate_hex(parameters["font_color"]):
                    run.font.color.rgb = RGBColor.from_string(parameters["font_color"].replace("#",""))
                if "font_size" in parameters:
                    run.font.size = Pt(parameters["font_size"])
                if "bold" in parameters:
                    run.font.bold = parameters["bold"]

    if "align_side" in parameters and "slide_height" in parameters and "slide_width" in parameters:
        align_side = parameters["align_side"]
        if "top" in align_side:
            shape.top = 0
        elif "bottom" in align_side:
            shape.top = parameters["slide_height"] - shape.height
        
        if "left" in align_side:
            shape.left = 0
        elif "right" in align_side:
            shape.right = parameters["slide_width"] - shape.width



def set_table_properties(shape, parameters):
    if not shape.has_table:
        return
    table = shape.table
    if "table_data" in parameters:
        table_data = parameters["table_data"]
        data_index = 0
        for row in table.rows:
            for col in row.cells:
                if data_index < len(table_data):
                    col.text = table_data[data_index]
                    data_index += 1
                else:
                    col.text = ""

    if "width" in parameters:
        col_width = toEmus(parameters["width"]/len(table.columns))
        for col in table.columns:
            col.width = col_width

    if "height" in parameters:
        row_height = toEmus(parameters["height"]/len(table.rows))
        for row in table.rows:
            row.height = row_height

    

def set_image_properties(picture, parameters):
    if "image_content" in parameters:
        im_bytes = generate_image(parameters["image_content"])
        im = Image.from_file(im_bytes)
        slide_part, rId = picture.part, picture._element.blip_rId
        image_part = slide_part.related_part(rId)
        image_part.blob = im._blob


#------------------------------------------------------------------------------------------------------------#


def modify_shape(ppt, slide_idx, input_parameters, model='gpt-4o-mini'):
    slide = ppt.slides[slide_idx]
    shape_idx = input_parameters["shape_index"]
    instructions = input_parameters["instructions"]
    shape_content = get_shape_content(slide, shape_idx)
    shape = slide.shapes[shape_idx]

    if 'PICTURE' in str(shape.shape_type):
        prompt = modify_picture_prompt
    elif 'CHART' in str(shape.shape_type):
        prompt = modify_chart_prompt
    elif 'TABLE' in str(shape.shape_type):
        prompt = modify_table_prompt
    else:
        prompt = modify_shape_prompt

    messages = [{"role":"system", "content":prompt}, {"role":"system", "content":shape_content}, {"role":"user", "content":instructions}]

    output_parameters = query(messages, json_mode=True, model=model)
    if "align_side" in output_parameters:
        output_parameters["slide_height"] = ppt.slide_height
        output_parameters["slide_width"] = ppt.slide_width
    set_image_properties(shape, output_parameters)
    set_table_properties(shape, output_parameters)
    set_shape_properties(shape, output_parameters)


    return f"Shape modified: {json.dumps(output_parameters)}"


def modify_background(ppt, slide_idx, parameters, model='gpt-4o-mini'):
    slide = ppt.slides[slide_idx]
    if "fill_color" in parameters:
        fill_color = parameters["fill_color"]

        if fill_color =="transparent":
            slide.background.fill.background()
        elif validate_hex(fill_color):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(fill_color.replace("#",""))

        return f"Background color modified: {fill_color}"


def insert_shape(ppt, slide_idx, parameters, model='gpt-4o-mini'):
    slide = ppt.slides[slide_idx]
    shape_type = parameters["shape_type"].strip().upper()
    if shape_type == "PICTURE":
        placeholder = BytesIO(open("src/data/placeholder.png", "rb").read())
        picture = slide.shapes.add_picture(placeholder, 0, 0)

    elif shape_type == "CHART":
        return "Chart operations not implemented"

    elif shape_type == "TABLE":
        messages = [{"role":"system", "content":select_table_dimensions_prompt}, {"role":"user", "content":parameters["instructions"]}]
        response = query(messages, json_mode=True, max_tokens=30, model=model)
        rows, cols = max(1, int(response["rows"])), max(1, int(response["columns"]))
        table = slide.shapes.add_table(rows, cols, 0, 0, toEmus(cols*40), toEmus(rows*12)).table
    elif shape_type == "TEXT_BOX":
        textbox = slide.shapes.add_textbox(0, 0, toEmus(50), toEmus(20))
    else:
        messages = [{"role":"system", "content":select_autoshape_prompt}, {"role":"system", "content":SHAPE_DICT}, {"role":"user", "content":parameters["instructions"]}]
        r = query(messages, json_mode=True, max_tokens=10, model=model)
        shape_id = int(r["id"])

        shape = slide.shapes.add_shape(shape_id, 0, 0, toEmus(10), toEmus(10))
        shape.fill.solid()

    response = f"Shape inserted: {shape_type} | "
    response += modify_shape(ppt, slide_idx, {"shape_index":len(slide.shapes)-1, "instructions":parameters["instructions"]})
    return response


def delete_shapes(ppt, slide_idx, parameters, model='gpt-4o-mini'):
    slide = ppt.slides[slide_idx]
    shape_indexes = parameters["shape_indexes"]
    shapes = slide.shapes
    shapes_to_remove = [shapes[idx].element for idx in shape_indexes]
    remove_count = 0
    for shape in shapes_to_remove:
        try: 
            shapes.element.remove(shape)
            remove_count += 1
        except:
            continue
    return f"{remove_count} shapes deleted"



def delete_all_shapes(ppt, slide_idx):
    slide = ppt.slides[slide_idx]
    shape_indexes = [i for i, _ in enumerate(slide.shapes)]
    response = f"Cleared slide {slide_idx} - " + delete_shapes(ppt, slide_idx, {"shape_indexes": shape_indexes})
    return response


def insert_slide(ppt, template_request, model='gpt-4o-mini'):
    layout_s = "Slide templates:\n" + "\n".join([f"{i} - {layout.name}" for i, layout in enumerate(ppt.slide_layouts)]) # Get all layouts given by slide master
    prompt = [{"role": "system", "content":select_layout_prompt},{"role": "system", "content":layout_s}, {"role":"user", "content":template_request}]

    r = query(prompt, json_mode=True, model=model, temperature=0, max_tokens=30)
    layout_idx = int(r["id"])

    slide_layout = ppt.slide_layouts[layout_idx]
    slide = ppt.slides.add_slide(slide_layout)

    response = f"Slide inserted with template: {slide_layout.name}"
    return response
