from src.agent import AgentPPT
from src.openai import generate_image, query
from src.ppt_reader import get_shape_content
from pptx import Presentation
import src.prompts as prompts
from src.ppt_writer import modify_shape, modify_background


def test_modify_shape():
    ppt = Presentation("test.pptx")
    slide = ppt.slides[0]
    instructions = "Change the text to 'Raccoons'"
    shape_idx = 0
    response = modify_shape(slide, {"instructions":instructions, "shape_index":shape_idx})
    
    shape_info = get_shape_content(slide, shape_idx, return_json=True)

    assert shape_info["text"] == "Raccoons"


def test_modify_picture():
    ppt = Presentation("test.pptx")
    slide = ppt.slides[0]
    instructions = "Change image to a large basketball Resize the image to a square."
    shape_idx = 3
    response = modify_shape(slide, {"instructions":instructions, "shape_index":shape_idx})

    shape_info = get_shape_content(slide, shape_idx, return_json=True)
    assert shape_info["size"]["height"] == shape_info["size"]["width"]